"""
build_adu_presentation.py
Generates an executive-grade 17-slide 16:9 PowerPoint (.pptx) and landscape PDF (.pdf)
portfolio covering all 5 S2A Modular Tiny Home & ADU models, CAD floorplans, interiors, and HaaS models.
"""

import os
import sys
import subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Paths
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DOWNLOADS_DIR = os.path.join(os.path.expanduser("~"), "Downloads")
ARTIFACT_DIR = "/Users/ericmiller/.gemini/antigravity-ide/brain/7edf9bcc-1ac5-41c6-befd-671b39214d4c"
APP_IMG_DIR = os.path.join(BASE_DIR, "haas-adu-configurator", "public", "images")

PPTX_OUTPUT_WORKSPACE = os.path.join(BASE_DIR, "S2A_Modular_ADU_HaaS_Portfolio.pptx")
PPTX_OUTPUT_DOWNLOADS = os.path.join(DOWNLOADS_DIR, "S2A_Modular_ADU_HaaS_Portfolio.pptx")
PPTX_OUTPUT_ARTIFACT = os.path.join(ARTIFACT_DIR, "S2A_Modular_ADU_HaaS_Portfolio.pptx")

PDF_OUTPUT_WORKSPACE = os.path.join(BASE_DIR, "S2A_Modular_ADU_HaaS_Portfolio.pdf")
PDF_OUTPUT_DOWNLOADS = os.path.join(DOWNLOADS_DIR, "S2A_Modular_ADU_HaaS_Portfolio.pdf")
PDF_OUTPUT_ARTIFACT = os.path.join(ARTIFACT_DIR, "S2A_Modular_ADU_HaaS_Portfolio.pdf")

# Images
HAVEN_IMG = os.path.join(APP_IMG_DIR, "haven.jpg")
HARMONY_IMG = os.path.join(APP_IMG_DIR, "harmony.jpg")
SIERRA_IMG = os.path.join(APP_IMG_DIR, "sierra.jpg")
MEADOW_IMG = os.path.join(APP_IMG_DIR, "meadow.jpg")
CASCADE_IMG = os.path.join(APP_IMG_DIR, "cascade.jpg")
INTERIOR_IMG = os.path.join(APP_IMG_DIR, "interior.jpg")
PATIO_IMG = os.path.join(APP_IMG_DIR, "patio.jpg")

# Floorplans & Blueprints
FP_HAVEN = os.path.join(APP_IMG_DIR, "floorplan_haven.png")
FP_HARMONY = os.path.join(APP_IMG_DIR, "floorplan_harmony.png")
FP_SIERRA = os.path.join(APP_IMG_DIR, "floorplan_sierra.png")
FP_MEADOW = os.path.join(APP_IMG_DIR, "floorplan_meadow.png")
FP_CASCADE = os.path.join(APP_IMG_DIR, "floorplan_cascade.png")

BP_HAVEN = os.path.join(APP_IMG_DIR, "blueprint_haven.png")
BP_HARMONY = os.path.join(APP_IMG_DIR, "blueprint_harmony.png")
BP_SIERRA = os.path.join(APP_IMG_DIR, "blueprint_sierra.png")
BP_MEADOW = os.path.join(APP_IMG_DIR, "blueprint_meadow.png")
BP_CASCADE = os.path.join(APP_IMG_DIR, "blueprint_cascade.png")

# Colors
C_DARK = RGBColor(15, 23, 42)       # #0f172a
C_NAVY = RGBColor(30, 41, 59)       # #1e293b
C_BLUE = RGBColor(37, 99, 235)      # #2563eb
C_LIGHT_BLUE = RGBColor(219, 234, 254) # #dbeafe
C_GREEN = RGBColor(22, 163, 74)     # #16a34a
C_WHITE = RGBColor(255, 255, 255)
C_GRAY = RGBColor(100, 116, 139)    # #64748b
C_LIGHT_BG = RGBColor(248, 250, 252)# #f8fafc
C_CARD_BG = RGBColor(241, 245, 249) # #f1f5f9
C_BORDER = RGBColor(203, 213, 225)

def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def set_bg(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, subtitle, title, light_theme=True):
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle.upper()
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = C_BLUE

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.5))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = C_DARK if light_theme else C_WHITE

    # SLIDE 1: COVER SLIDE
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1, C_DARK)
    
    if os.path.exists(HAVEN_IMG):
        s1.shapes.add_picture(HAVEN_IMG, Inches(6.8), Inches(0.8), Inches(5.8), Inches(5.9))

    cover_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.5))
    tf1 = cover_box.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "S2A MODULAR • HAAS PORTFOLIO"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_BLUE
    
    p2 = tf1.add_paragraph()
    p2.text = "Next-Gen Modular ADU &\nTiny Home Collection"
    p2.font.size = Pt(30)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.space_before = Pt(12)

    p3 = tf1.add_paragraph()
    p3.text = "Complete Architectural CAD Floorplans & Turnkey Housing-as-a-Service (HaaS) Revenue Suite."
    p3.font.size = Pt(13)
    p3.font.color.rgb = C_LIGHT_BLUE
    p3.space_before = Pt(14)

    p4 = tf1.add_paragraph()
    p4.text = "Zero Upfront Capital • Passive Homeowner Revenue • Turnkey Operator Delivery"
    p4.font.size = Pt(11)
    p4.font.color.rgb = C_GREEN
    p4.space_before = Pt(16)

    # SLIDE 2: THE HAAS REVENUE-SHARE MODEL
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2, C_LIGHT_BG)
    add_header(s2, "Housing-as-a-Service (HaaS) Economics", "The Zero-CapEx Turnkey Backyard Revenue Model")

    cards_data = [
        ("1. Homeowner Benefit", "$500 – $1,000 / Mo Net", "Homeowners license unused backyard space for zero upfront cost. Operator funds permitting, foundation, unit fabrication, crane installation, and tenant management.", C_BLUE),
        ("2. Operator Economics", "$1,500 – $2,200 / Mo Net Yield", "High-ADR mid-term corporate executive, traveling medical, and high-quality long-term tenant placement yields 18–24 month capital payback per installation.", C_GREEN),
        ("3. Precision Modular Speed", "30-Day Delivery Timeline", "S2A factory-built offsite construction eliminates 90% of neighborhood disruption. Site utility trenching and crane set takes under 48 hours.", C_DARK)
    ]

    for i, (head, val, desc, color) in enumerate(cards_data):
        left = Inches(0.8 + i * 3.95)
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(3.8), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = C_BORDER

        tb = s2.shapes.add_textbox(left + Inches(0.2), Inches(1.7), Inches(3.4), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = head
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_DARK

        p_val = tf.add_paragraph()
        p_val.text = val
        p_val.font.size = Pt(16)
        p_val.font.bold = True
        p_val.font.color.rgb = color
        p_val.space_before = Pt(8)

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = C_GRAY
        p_desc.space_before = Pt(12)

    # FUNCTION TO ADD MODEL OVERVIEW SLIDE
    def add_model_slide(img_path, subtitle, title, dim, area, rent, split, features, ideal_for):
        s = prs.slides.add_slide(blank_layout)
        set_bg(s, C_LIGHT_BG)
        add_header(s, subtitle, title)

        if os.path.exists(img_path):
            s.shapes.add_picture(img_path, Inches(0.8), Inches(1.4), Inches(6.2), Inches(5.4))

        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.4), Inches(5.3), Inches(5.4))
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = C_BORDER

        tb = s.shapes.add_textbox(Inches(7.4), Inches(1.5), Inches(4.9), Inches(5.2))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"DIMENSIONS: {dim}  |  LIVING AREA: {area}"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = C_BLUE

        p_rent = tf.add_paragraph()
        p_rent.text = f"Est. Monthly Rent: {rent}"
        p_rent.font.size = Pt(16)
        p_rent.font.bold = True
        p_rent.font.color.rgb = C_GREEN
        p_rent.space_before = Pt(4)

        p_split = tf.add_paragraph()
        p_split.text = f"Homeowner Passive Net: {split}"
        p_split.font.size = Pt(13)
        p_split.font.bold = True
        p_split.font.color.rgb = C_DARK

        p_h = tf.add_paragraph()
        p_h.text = "Architectural & Blueprint Highlights:"
        p_h.font.size = Pt(11.5)
        p_h.font.bold = True
        p_h.font.color.rgb = C_NAVY
        p_h.space_before = Pt(10)

        for feat in features:
            pf = tf.add_paragraph()
            pf.text = f"• {feat}"
            pf.font.size = Pt(10)
            pf.font.color.rgb = C_GRAY
            pf.space_before = Pt(3)

        p_ideal = tf.add_paragraph()
        p_ideal.text = f"Ideal Placement: {ideal_for}"
        p_ideal.font.size = Pt(10)
        p_ideal.font.bold = True
        p_ideal.font.color.rgb = C_BLUE
        p_ideal.space_before = Pt(10)

    # FUNCTION TO ADD FLOORPLAN SLIDE
    def add_floorplan_slide(fp_path, bp_path, subtitle, title, room_specs):
        s = prs.slides.add_slide(blank_layout)
        set_bg(s, C_LIGHT_BG)
        add_header(s, subtitle, title)

        # Left Floorplan Image Box
        if os.path.exists(fp_path):
            card_fp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(7.0), Inches(5.4))
            card_fp.fill.solid()
            card_fp.fill.fore_color.rgb = C_WHITE
            card_fp.line.color.rgb = C_BORDER
            s.shapes.add_picture(fp_path, Inches(1.0), Inches(1.6), Inches(6.6), Inches(5.0))

        # Right Specs Card
        card_spec = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.4), Inches(4.5), Inches(5.4))
        card_spec.fill.solid()
        card_spec.fill.fore_color.rgb = C_WHITE
        card_spec.line.color.rgb = C_BORDER

        tb = s.shapes.add_textbox(Inches(8.2), Inches(1.6), Inches(4.1), Inches(5.0))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "OFFICIAL CAD BLUEPRINT SPECIFICATIONS"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_BLUE

        p_sub = tf.add_paragraph()
        p_sub.text = "Factory Framed & State Approved"
        p_sub.font.size = Pt(14)
        p_sub.font.bold = True
        p_sub.font.color.rgb = C_DARK
        p_sub.space_before = Pt(4)

        for room, detail in room_specs:
            p_rm = tf.add_paragraph()
            p_rm.text = f"• {room}: {detail}"
            p_rm.font.size = Pt(10)
            p_rm.font.color.rgb = C_GRAY
            p_rm.space_before = Pt(6)

        p_stamp = tf.add_paragraph()
        p_stamp.text = "✓ Sealed Architectural Drawing Number A1"
        p_stamp.font.size = Pt(10)
        p_stamp.font.bold = True
        p_stamp.font.color.rgb = C_GREEN
        p_stamp.space_before = Pt(14)

    # 1. HAVEN
    add_model_slide(
        HAVEN_IMG,
        "Model 01 • Modern Luxury ADU",
        "The Haven (14' × 29' • 350 Sq. Ft.)",
        "14' × 29'-0\"", "350 Sq. Ft.",
        "$2,200 – $2,950 / mo", "$650 – $950 / mo",
        [
            "Full rooftop sky lounge deck with spiral staircase access",
            "4'-0\" integrated covered front porch",
            "Open kitchenette & living room with roof skylight access",
            "Private bedroom with double wardrobe closet & full bathroom",
            "Large sliding glass patio doors for indoor-outdoor flow"
        ],
        "Upscale suburban backyards, premium view lots, executive retreats."
    )
    add_floorplan_slide(
        FP_HAVEN, BP_HAVEN,
        "Model 01 • Architectural Plan",
        "The Haven Floorplan (14' × 29' • 350 Sq. Ft.)",
        [
            ("Covered Porch", "4'-0\" integrated covered front entry porch with recessed downlights."),
            ("Great Room & Kitchen", "Open 14' wide living space with roof skylight access hatch."),
            ("Private Bedroom", "Ground floor bedroom with dual 2442SH egress windows & double closet."),
            ("Full Bathroom", "3-3060 shower, vanity, toilet, and dedicated stack washer/dryer."),
            ("Rooftop Deck", "Optional rooftop terrace accessed via site-installed spiral stairs.")
        ]
    )

    # 2. HARMONY
    add_model_slide(
        HARMONY_IMG,
        "Model 02 • Modern Craftsman Cottage",
        "The Harmony (14' × 24' • 348 Sq. Ft.)",
        "14' × 24'-0\"", "348 Sq. Ft.",
        "$1,850 – $2,350 / mo", "$550 – $800 / mo",
        [
            "Classic pitched gabled roofline with architectural shingles",
            "Welcoming covered front portico with artisan coach lights",
            "Smart-lap horizontal white siding with black frame windows",
            "Spacious Bed/Living open studio layout with full kitchen",
            "Complete bathroom with 3-3060 shower and stack W/D"
        ],
        "Traditional HOA neighborhoods, in-law suites, historic districts."
    )
    add_floorplan_slide(
        FP_HARMONY, BP_HARMONY,
        "Model 02 • Architectural Plan",
        "The Harmony Floorplan (14' × 24' • 348 Sq. Ft.)",
        [
            ("Bed/Living Suite", "Large 14' × 12' open living/bedroom area with 3060SH windows."),
            ("Full Kitchen", "Complete cooktop, sink, refrigerator, and pantry cabinet."),
            ("Full Bathroom", "Enclosed private bath with 3-3060 fiberglass shower stall."),
            ("Laundry & Storage", "Built-in stackable W/D closet with 2/4 INT door."),
            ("Front Portico", "Covered stoop with classic artisan architectural trim.")
        ]
    )

    # 3. SIERRA
    add_model_slide(
        SIERRA_IMG,
        "Model 03 • Mono-Pitch Shed Loft",
        "The Sierra (14' × 31'-6\" • 420 Sq. Ft.)",
        "14' × 31'-6\"", "420 Sq. Ft.",
        "$2,400 – $3,100 / mo", "$600 – $900 / mo",
        [
            "Modern mono-pitch shed roofline with clerestory transoms",
            "Vaulted Great Room with mezzanine sleeping loft above",
            "Ground floor private bedroom with large closet & egress",
            "Board-and-batten dark charcoal siding with natural oak inlays",
            "Front and rear entry decks with step lighting"
        ],
        "Medium-to-large residential lots, remote tech professionals."
    )
    add_floorplan_slide(
        FP_SIERRA, BP_SIERRA,
        "Model 03 • Architectural Plan",
        "The Sierra Floorplan (14' × 31'-6\" • 420 Sq. Ft.)",
        [
            ("Vaulted Great Room", "Expansive 30'-0\" span with double front sliding glass doors."),
            ("Mezzanine Loft", "Upper sleeping/storage loft open to below with access ladder."),
            ("Private Bedroom", "Ground floor rear bedroom with full closet and egress windows."),
            ("Full Kitchen & Dining", "Cooktop, prep island, sink, and high clerestory windows."),
            ("Outdoor Stoops", "Dual front and rear site-built wooden entry stoops.")
        ]
    )

    # 4. MEADOW
    add_model_slide(
        MEADOW_IMG,
        "Model 04 • Flagship Rooftop Terrace",
        "The Meadow (15' × 31'-9\" • 435 Sq. Ft.)",
        "15' × 31'-9\"", "435 Sq. Ft.",
        "$2,600 – $3,400 / mo", "$700 – $1,000 / mo",
        [
            "Full-footprint rooftop sun terrace with slatted privacy railings",
            "Exterior open staircase connecting yard directly to rooftop",
            "Symmetrical front facade with flanked oversized black windows",
            "Dual lofts (sleeping + storage/lounge) and plant ledge top",
            "Full gourmet kitchen with dining peninsula and rear patio door"
        ],
        "Wide suburban parcels (50'+ width), high-ADR luxury travel rentals."
    )
    add_floorplan_slide(
        FP_MEADOW, BP_MEADOW,
        "Model 04 • Architectural Plan",
        "The Meadow Floorplan (15' × 31'-9\" • 435 Sq. Ft.)",
        [
            ("Dual Lofts", "Dual mezzanine lofts (Loft 1 + Loft 2) with plant ledge top."),
            ("Exterior Staircase", "2'-9\" exterior open stair ascending to full rooftop terrace."),
            ("Gourmet Kitchen", "Central kitchen with peninsula dining & 4836SL windows."),
            ("Private Suite", "Private master bedroom with closet and stack W/D closet."),
            ("Rear Patio Egress", "6/0 Patio sliding door opening directly onto garden deck.")
        ]
    )

    # 5. CASCADE
    add_model_slide(
        CASCADE_IMG,
        "Model 05 • Ultra-Compact High-Yield Infill",
        "The Cascade (12' × 27'-9\" • 300 Sq. Ft.)",
        "12' × 27'-9\"", "300 Sq. Ft.",
        "$1,650 – $2,100 / mo", "$500 – $750 / mo",
        [
            "Slimline 12-foot exterior profile engineered for narrow lots",
            "Industrial steel ship's ladder ascending to observation deck",
            "Efficient open Living/Dining, Kitchenette, and Bath",
            "Ground floor bedroom with upper loft overhead",
            "Fits 95% of standard residential parcels without setback variance"
        ],
        "Narrow side yards, urban infill parcels, maximum ROI density."
    )
    add_floorplan_slide(
        FP_CASCADE, BP_CASCADE,
        "Model 05 • Architectural Plan",
        "The Cascade Floorplan (12' × 27'-9\" • 300 Sq. Ft.)",
        [
            ("12' Slimline Footprint", "26'-0\" main box length engineered for tight setbacks."),
            ("Ship's Ladder Access", "2'-9\" side ship's ladder leading to rooftop observation deck."),
            ("Living/Dining & Kitchenette", "Efficient open plan with 4836SL window and linen storage."),
            ("Bedroom & Loft", "Private sleeping quarters with bonus overhead storage loft."),
            ("Rear Patio Door", "6/0 Patio sliding door for indoor-outdoor access.")
        ]
    )

    # SLIDE 13: INTERIOR ARCHITECTURE
    s13 = prs.slides.add_slide(blank_layout)
    set_bg(s13, C_LIGHT_BG)
    add_header(s13, "Interior Architecture & Finish Standards", "Scandinavian Modern Luxury & Space Optimization")

    if os.path.exists(INTERIOR_IMG):
        s13.shapes.add_picture(INTERIOR_IMG, Inches(0.8), Inches(1.4), Inches(6.2), Inches(5.4))

    card13 = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.4), Inches(5.3), Inches(5.4))
    card13.fill.solid()
    card13.fill.fore_color.rgb = C_WHITE
    card13.line.color.rgb = C_BORDER

    tb13 = s13.shapes.add_textbox(Inches(7.4), Inches(1.5), Inches(4.9), Inches(5.2))
    tf13 = tb13.text_frame
    tf13.word_wrap = True

    p = tf13.paragraphs[0]
    p.text = "PREMIUM FACTORY-FINISHED INTERIORS"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_BLUE

    points = [
        ("Vaulted Timber Ceilings", "Exposed natural wood trusses create expansive spatial volume."),
        ("Mezzanine Sleeping Lofts", "Adds 100+ sq ft of functional bonus space for sleeping or storage."),
        ("Quartz Waterfall Kitchens", "Durable quartz countertops, custom oak flat-panel cabinetry."),
        ("European Oak Hardwood", "Waterproof luxury engineered plank flooring throughout."),
        ("Designer Bathrooms", "Walk-in glass showers, floating vanities, matte black hardware."),
        ("Energy Efficiency", "All-electric mini-split HVAC, smart thermostats, tankless hot water.")
    ]

    for title_pt, desc_pt in points:
        pt = tf13.add_paragraph()
        pt.text = f"• {title_pt}: {desc_pt}"
        pt.font.size = Pt(10)
        pt.font.color.rgb = C_DARK
        pt.space_before = Pt(5)

    # SLIDE 14: REAR PATIO & OUTDOOR LIVING
    s14 = prs.slides.add_slide(blank_layout)
    set_bg(s14, C_LIGHT_BG)
    add_header(s14, "Seamless Indoor-Outdoor Living", "Private Backyard Entertaining Patios & Sun Decks")

    if os.path.exists(PATIO_IMG):
        s14.shapes.add_picture(PATIO_IMG, Inches(0.8), Inches(1.4), Inches(6.2), Inches(5.4))

    card14 = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.4), Inches(5.3), Inches(5.4))
    card14.fill.solid()
    card14.fill.fore_color.rgb = C_WHITE
    card14.line.color.rgb = C_BORDER

    tb14 = s14.shapes.add_textbox(Inches(7.4), Inches(1.5), Inches(4.9), Inches(5.2))
    tf14 = tb14.text_frame
    tf14.word_wrap = True

    p = tf14.paragraphs[0]
    p.text = "ELEVATING THE BACKYARD EXPERIENCE"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_BLUE

    points9 = [
        ("Floor-to-Ceiling Sliders", "Seamless transition between interior living room and private garden deck."),
        ("Dedicated Outdoor Dining", "Space for outdoor tables, lounge chairs, and BBQ entertaining."),
        ("Integrated Step Lighting", "Warm low-voltage LED pathway lighting for evening safety and ambiance."),
        ("Privacy Screening", "Architectural slat walls protect homeowner and tenant privacy."),
        ("Minimal Yard Footprint", "Leaves 60-70% of standard backyards open for landscaping and lawns.")
    ]

    for title_pt, desc_pt in points9:
        pt = tf14.add_paragraph()
        pt.text = f"• {title_pt}: {desc_pt}"
        pt.font.size = Pt(10)
        pt.font.color.rgb = C_DARK
        pt.space_before = Pt(6)

    # SLIDE 15: COMPARISON MATRIX
    s15 = prs.slides.add_slide(blank_layout)
    set_bg(s15, C_LIGHT_BG)
    add_header(s15, "Portfolio Comparison", "Comprehensive S2A Modular ADU Lineup")

    rows, cols = 6, 6
    left, top, width, height = Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2)
    table_shape = s15.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table

    headers = ["Model", "Dimensions", "Living Area", "Key Blueprint Feature", "Est. Rent", "Homeowner Split"]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_WHITE

    matrix_data = [
        ("Cascade", "12' × 27'-9\"", "300 sq. ft.", "12' Slim Infill • Rooftop Ladder", "$1,650 – $2,100", "$500 – $750 / mo"),
        ("Harmony", "14' × 24'-0\"", "348 sq. ft.", "Pitched Roof • Covered Portico", "$1,850 – $2,350", "$550 – $800 / mo"),
        ("Haven", "14' × 29'-0\"", "350 sq. ft.", "4' Porch • Rooftop Sky Lounge", "$2,200 – $2,950", "$650 – $950 / mo"),
        ("Sierra", "14' × 31'-6\"", "420 sq. ft.", "Mono-Pitch Shed • Clerestory Loft", "$2,400 – $3,100", "$600 – $900 / mo"),
        ("Meadow", "15' × 31'-9\"", "435 sq. ft.", "Dual Lofts • Full Rooftop Terrace", "$2,600 – $3,400", "$700 – $1,000 / mo"),
    ]

    for i, row in enumerate(matrix_data):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_WHITE if i % 2 == 0 else C_CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = C_DARK if j != 0 else C_BLUE
            if j == 0 or j >= 4:
                p.font.bold = True

    # SLIDE 16: DEPLOYMENT TIMELINE
    s16 = prs.slides.add_slide(blank_layout)
    set_bg(s16, C_LIGHT_BG)
    add_header(s16, "Turnkey Delivery Timeline", "From Site Audit to Monthly Cash Flow in 30–45 Days")

    timeline = [
        ("Day 1–5", "Site Feasibility & Digital Permit", "GIS setback validation, utility hookup mapping, pre-approved architectural plan submission."),
        ("Day 6–20", "Offsite Factory Modular Build", "Precision factory framing, plumbing, electrical, quartz counters, and paint under climate-controlled conditions."),
        ("Day 21–25", "Foundation & Trenching", "Helical pile / concrete pad foundation and underground sewer/power/water utility trenching in yard."),
        ("Day 26–30", "Crane Delivery & Installation", "Single-day crane lift over main house roofline, utility tie-in, and final city inspection sign-off."),
        ("Day 31+", "Tenant Placement & Revenue", "Furnishing, professional photography, tenant onboarding, and automated monthly revenue-share payouts.")
    ]

    for i, (day, title_tl, desc_tl) in enumerate(timeline):
        card = s16.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5 + i * 1.05), Inches(11.7), Inches(0.9))
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = C_BORDER

        tb = s16.shapes.add_textbox(Inches(1.0), Inches(1.55 + i * 1.05), Inches(11.3), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{day}: {title_tl}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_BLUE

        p_desc = tf.add_paragraph()
        p_desc.text = desc_tl
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = C_GRAY

    # SLIDE 17: CALL TO ACTION
    s17 = prs.slides.add_slide(blank_layout)
    set_bg(s17, C_DARK)

    tb17 = s17.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.5))
    tf17 = tb17.text_frame
    tf17.word_wrap = True

    p = tf17.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "UNLOCK YOUR BACKYARD'S PASSIVE INCOME"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_BLUE

    p2 = tf17.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "Get Your Free 48-Hour Site Feasibility &\nRevenue-Share Audit"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.space_before = Pt(12)

    p3 = tf17.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.text = "We assess your parcel geometry, local ADU zoning, utility trenching distances,\nand provide a guaranteed monthly passive income estimate."
    p3.font.size = Pt(13)
    p3.font.color.rgb = C_LIGHT_BLUE
    p3.space_before = Pt(14)

    p4 = tf17.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    p4.text = "S2A Modular • Housing-as-a-Service (HaaS) Solutions"
    p4.font.size = Pt(12)
    p4.font.bold = True
    p4.font.color.rgb = C_GREEN
    p4.space_before = Pt(20)

    prs.save(PPTX_OUTPUT_WORKSPACE)
    prs.save(PPTX_OUTPUT_DOWNLOADS)
    prs.save(PPTX_OUTPUT_ARTIFACT)
    print(f"PPTX successfully created (17 slides):\n- Downloads: {PPTX_OUTPUT_DOWNLOADS}\n- Workspace: {PPTX_OUTPUT_WORKSPACE}")

def create_pdf():
    html_file = os.path.join(BASE_DIR, "cache", "deck_temp.html")
    
    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<title>S2A Modular ADU & Tiny Home Collection - Architectural CAD & HaaS Presentation</title>
<style>
  @page {{
    size: 11in 6.1875in;
    margin: 0;
  }}
  * {{
    box-sizing: border-box;
    -webkit-print-color-adjust: exact !important;
    print-color-adjust: exact !important;
  }}
  body {{
    font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, Helvetica, Arial, sans-serif;
    margin: 0;
    padding: 0;
    background: #f8fafc;
    color: #1e293b;
    font-size: 11px;
  }}
  .slide {{
    width: 11in;
    height: 6.1875in;
    padding: 0.35in 0.55in;
    position: relative;
    page-break-after: always;
    break-after: page;
    background: #f8fafc;
    display: flex;
    flex-direction: column;
  }}
  .slide-dark {{
    background: #0f172a;
    color: #ffffff;
  }}
  .slide-header {{
    margin-bottom: 12px;
  }}
  .sub-tag {{
    font-size: 8.5px;
    font-weight: 800;
    letter-spacing: 1.5px;
    text-transform: uppercase;
    color: #2563eb;
    margin-bottom: 2px;
  }}
  h2 {{
    font-size: 18px;
    font-weight: 800;
    margin: 0;
    color: #0f172a;
    letter-spacing: -0.4px;
  }}
  .slide-dark h2 {{
    color: #ffffff;
  }}
  .two-col {{
    display: grid;
    grid-template-columns: 1.15fr 1fr;
    gap: 16px;
    flex: 1;
    align-items: stretch;
  }}
  .model-img {{
    width: 100%;
    height: 100%;
    max-height: 4.5in;
    object-fit: cover;
    border-radius: 6px;
    border: 1px solid #cbd5e1;
  }}
  .fp-img-box {{
    background: #ffffff;
    border: 1px solid #cbd5e1;
    border-radius: 6px;
    padding: 12px;
    display: flex;
    align-items: center;
    justify-content: center;
    height: 100%;
    max-height: 4.5in;
  }}
  .fp-img {{
    width: 100%;
    height: 100%;
    object-fit: contain;
  }}
  .info-card {{
    background: #ffffff;
    border: 1px solid #e2e8f0;
    border-radius: 6px;
    padding: 12px 14px;
    display: flex;
    flex-direction: column;
    justify-content: space-between;
  }}
  .badge-dim {{
    font-size: 9px;
    font-weight: 700;
    color: #2563eb;
    background: #eff6ff;
    padding: 2px 6px;
    border-radius: 4px;
    display: inline-block;
    margin-bottom: 6px;
  }}
  .rent-val {{
    font-size: 15px;
    font-weight: 800;
    color: #16a34a;
  }}
  .split-val {{
    font-size: 12px;
    font-weight: 700;
    color: #0f172a;
    margin-bottom: 8px;
  }}
  .bullet-list {{
    margin: 0;
    padding-left: 14px;
    font-size: 9.5px;
    color: #475569;
    line-height: 1.45;
  }}
  .bullet-list li {{
    margin-bottom: 3px;
  }}
  .ideal-tag {{
    font-size: 8.5px;
    font-weight: 700;
    color: #2563eb;
    margin-top: 6px;
  }}
  table {{
    width: 100%;
    border-collapse: collapse;
    font-size: 9.5px;
    margin-top: 10px;
    background: #ffffff;
    border-radius: 6px;
    overflow: hidden;
    border: 1px solid #e2e8f0;
  }}
  th {{
    background: #0f172a;
    color: #ffffff;
    text-align: left;
    padding: 7px 10px;
    font-weight: 700;
  }}
  td {{
    padding: 6px 10px;
    border-bottom: 1px solid #e2e8f0;
  }}
  tr:nth-child(even) td {{
    background: #f8fafc;
  }}
  .footer-bar {{
    position: absolute;
    bottom: 0.2in;
    left: 0.55in;
    right: 0.55in;
    display: flex;
    justify-content: space-between;
    font-size: 7.5px;
    color: #94a3b8;
    border-top: 1px solid #e2e8f0;
    padding-top: 4px;
  }}
  .slide-dark .footer-bar {{
    border-top-color: #334155;
    color: #64748b;
  }}
</style>
</head>
<body>

  <!-- SLIDE 1: COVER -->
  <div class="slide slide-dark" style="justify-content: center; padding-left: 0.8in;">
    <div style="display: grid; grid-template-columns: 1.1fr 1fr; gap: 20px; align-items: center; height: 100%;">
      <div>
        <div class="sub-tag">S2A MODULAR &bull; HAAS PORTFOLIO</div>
        <h1 style="font-size: 28px; font-weight: 800; margin: 8px 0; color: #ffffff; line-height: 1.15;">
          Next-Gen Modular ADU &<br/>Tiny Home Collection
        </h1>
        <p style="font-size: 11.5px; color: #93c5fd; margin: 10px 0 16px 0; line-height: 1.4;">
          Architectural CAD Floorplans & Turnkey Housing-as-a-Service (HaaS) Revenue Suite: Haven, Harmony, Sierra, Meadow & Cascade.
        </p>
        <div style="font-size: 10px; font-weight: 700; color: #4ade80;">
          &bull; Zero Upfront Capital &bull; Passive Homeowner Revenue &bull; Official Sealed Blueprints
        </div>
      </div>
      <div>
        <img src="file://{HAVEN_IMG}" style="width: 100%; height: 4.6in; object-fit: cover; border-radius: 6px; border: 1px solid #334155;" />
      </div>
    </div>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 1 of 17</span>
    </div>
  </div>

  <!-- SLIDE 2: HAAS ECONOMICS -->
  <div class="slide">
    <div class="slide-header">
      <div class="sub-tag">HOUSING-AS-A-SERVICE (HAAS) ECONOMICS</div>
      <h2>The Zero-CapEx Turnkey Backyard Revenue Model</h2>
    </div>
    <div style="display: grid; grid-template-columns: 1fr 1fr 1fr; gap: 14px; flex: 1;">
      <div class="info-card">
        <div>
          <div style="font-size: 12px; font-weight: 800; color: #0f172a; margin-bottom: 4px;">1. Homeowner Benefit</div>
          <div style="font-size: 14px; font-weight: 800; color: #2563eb; margin-bottom: 8px;">$500 &ndash; $1,000 / Mo Net</div>
          <p style="font-size: 9.5px; color: #64748b; line-height: 1.45;">
            Homeowners license unused backyard space for zero upfront cost. Operator funds permitting, foundation, unit fabrication, crane installation, and tenant management.
          </p>
        </div>
        <div style="font-size: 8px; font-weight: 700; color: #16a34a;">100% Passive &bull; Zero Landlord Headaches</div>
      </div>
      <div class="info-card">
        <div>
          <div style="font-size: 12px; font-weight: 800; color: #0f172a; margin-bottom: 4px;">2. Operator Economics</div>
          <div style="font-size: 14px; font-weight: 800; color: #16a34a; margin-bottom: 8px;">$1,500 &ndash; $2,200 / Mo Net Yield</div>
          <p style="font-size: 9.5px; color: #64748b; line-height: 1.45;">
            High-ADR mid-term corporate executive, traveling medical, and high-quality long-term tenant placement yields 18–24 month capital payback per installation.
          </p>
        </div>
        <div style="font-size: 8px; font-weight: 700; color: #2563eb;">Scalable Distributed Real Estate Portfolio</div>
      </div>
      <div class="info-card">
        <div>
          <div style="font-size: 12px; font-weight: 800; color: #0f172a; margin-bottom: 4px;">3. Precision Modular Speed</div>
          <div style="font-size: 14px; font-weight: 800; color: #0f172a; margin-bottom: 8px;">30-Day Delivery Timeline</div>
          <p style="font-size: 9.5px; color: #64748b; line-height: 1.45;">
            S2A factory-built offsite construction eliminates 90% of neighborhood disruption. Site utility trenching and crane set takes under 48 hours.
          </p>
        </div>
        <div style="font-size: 8px; font-weight: 700; color: #0f172a;">Factory Quality &bull; Pre-Approved Plans</div>
      </div>
    </div>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 2 of 17</span>
    </div>
  </div>

  <!-- SLIDE 3: HAVEN 3D -->
  <div class="slide">
    <div class="slide-header">
      <div class="sub-tag">MODEL 01 &bull; 3D EXTERIOR & REVENUE</div>
      <h2>The Haven &bull; 14' &times; 29' (350 Sq. Ft.)</h2>
    </div>
    <div class="two-col">
      <img src="file://{HAVEN_IMG}" class="model-img" />
      <div class="info-card">
        <div>
          <div class="badge-dim">DIMENSIONS: 14' &times; 29'-0" &bull; 350 SQ. FT.</div>
          <div class="rent-val">Est. Monthly Rent: $2,200 &ndash; $2,950 / mo</div>
          <div class="split-val">Homeowner Passive Net: $650 &ndash; $950 / mo</div>
          <div style="font-weight: 700; font-size: 10px; margin-bottom: 4px;">Architectural Highlights:</div>
          <ul class="bullet-list">
            <li>Full rooftop sky lounge deck with spiral staircase access</li>
            <li>4'-0" integrated covered front porch with exterior lighting</li>
            <li>Open kitchenette & living room with roof skylight access</li>
            <li>Private bedroom with double wardrobe closet & full bath</li>
            <li>Large sliding glass patio doors for indoor-outdoor flow</li>
          </ul>
        </div>
        <div class="ideal-tag">Ideal Placement: Upscale suburban backyards, premium view lots, executive retreats.</div>
      </div>
    </div>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 3 of 17</span>
    </div>
  </div>

  <!-- SLIDE 4: HAVEN FLOORPLAN -->
  <div class="slide">
    <div class="slide-header">
      <div class="sub-tag">MODEL 01 &bull; ARCHITECTURAL CAD FLOORPLAN</div>
      <h2>The Haven CAD Blueprint (14' &times; 29'-0" &bull; 350 Sq. Ft.)</h2>
    </div>
    <div class="two-col">
      <div class="fp-img-box">
        <img src="file://{FP_HAVEN}" class="fp-img" />
      </div>
      <div class="info-card">
        <div>
          <div class="badge-dim">SEALED DRAWING NUMBER: A1 &bull; S2A MODULAR</div>
          <div style="font-weight: 800; font-size: 12px; color: #0f172a; margin-bottom: 6px;">Detailed Room Layout:</div>
          <ul class="bullet-list">
            <li><strong>4'-0" Covered Porch:</strong> Integrated entry porch with recessed soffit lighting.</li>
            <li><strong>Open Great Room:</strong> Kitchenette and living area with roof skylight access hatch.</li>
            <li><strong>Private Master Bedroom:</strong> Ground floor with dual 2442SH windows and double wardrobe.</li>
            <li><strong>Designer Full Bath:</strong> 3-3060 fiberglass shower, vanity, toilet, and stack W/D closet.</li>
            <li><strong>Rooftop Access:</strong> Optional spiral stairs leading to full-length observation deck.</li>
          </ul>
        </div>
        <div style="font-size: 9px; font-weight: 700; color: #16a34a;">✓ Pre-Approved California SB 9 / Statewide ADU Building Code Compliant</div>
      </div>
    </div>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 4 of 17</span>
    </div>
  </div>

  <!-- SLIDE 5: HARMONY 3D -->
  <div class="slide">
    <div class="slide-header">
      <div class="sub-tag">MODEL 02 &bull; 3D EXTERIOR & REVENUE</div>
      <h2>The Harmony &bull; 14' &times; 24' (348 Sq. Ft.)</h2>
    </div>
    <div class="two-col">
      <img src="file://{HARMONY_IMG}" class="model-img" />
      <div class="info-card">
        <div>
          <div class="badge-dim">DIMENSIONS: 14' &times; 24'-0" &bull; 348 SQ. FT.</div>
          <div class="rent-val">Est. Monthly Rent: $1,850 &ndash; $2,350 / mo</div>
          <div class="split-val">Homeowner Passive Net: $550 &ndash; $800 / mo</div>
          <div style="font-weight: 700; font-size: 10px; margin-bottom: 4px;">Architectural Highlights:</div>
          <ul class="bullet-list">
            <li>Classic pitched gabled roofline with architectural asphalt shingles</li>
            <li>Welcoming covered front portico with artisan coach lights</li>
            <li>Smart-lap horizontal white siding with black frame windows</li>
            <li>Spacious Bed/Living open studio layout with full kitchen</li>
            <li>Complete bathroom with 3-3060 shower and stack W/D</li>
          </ul>
        </div>
        <div class="ideal-tag">Ideal Placement: Traditional HOA neighborhoods, in-law suites, historic districts.</div>
      </div>
    </div>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 5 of 17</span>
    </div>
  </div>

  <!-- SLIDE 6: HARMONY FLOORPLAN -->
  <div class="slide">
    <div class="slide-header">
      <div class="sub-tag">MODEL 02 &bull; ARCHITECTURAL CAD FLOORPLAN</div>
      <h2>The Harmony CAD Blueprint (14' &times; 24'-0" &bull; 348 Sq. Ft.)</h2>
    </div>
    <div class="two-col">
      <div class="fp-img-box">
        <img src="file://{FP_HARMONY}" class="fp-img" />
      </div>
      <div class="info-card">
        <div>
          <div class="badge-dim">SEALED DRAWING NUMBER: A1 &bull; S2A MODULAR</div>
          <div style="font-weight: 800; font-size: 12px; color: #0f172a; margin-bottom: 6px;">Detailed Room Layout:</div>
          <ul class="bullet-list">
            <li><strong>Bed/Living Open Studio:</strong> Generous 14' wide studio living with 3060SH windows.</li>
            <li><strong>Full Chef's Kitchen:</strong> Range, exhaust, stainless sink, pantry, and refrigerator.</li>
            <li><strong>Enclosed Bathroom:</strong> 3-3060 shower stall, sink vanity, and privacy door.</li>
            <li><strong>Stack W/D Enclosure:</strong> Dedicated laundry utility space.</li>
            <li><strong>Entry Stoop:</strong> Front portico with craftsman columns.</li>
          </ul>
        </div>
        <div style="font-size: 9px; font-weight: 700; color: #16a34a;">✓ Pre-Approved California SB 9 / Statewide ADU Building Code Compliant</div>
      </div>
    </div>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 6 of 17</span>
    </div>
  </div>

  <!-- SLIDE 7: SIERRA 3D -->
  <div class="slide">
    <div class="slide-header">
      <div class="sub-tag">MODEL 03 &bull; 3D EXTERIOR & REVENUE</div>
      <h2>The Sierra &bull; 14' &times; 31'-6" (420 Sq. Ft.)</h2>
    </div>
    <div class="two-col">
      <img src="file://{SIERRA_IMG}" class="model-img" />
      <div class="info-card">
        <div>
          <div class="badge-dim">DIMENSIONS: 14' &times; 31'-6" &bull; 420 SQ. FT.</div>
          <div class="rent-val">Est. Monthly Rent: $2,400 &ndash; $3,100 / mo</div>
          <div class="split-val">Homeowner Passive Net: $600 &ndash; $900 / mo</div>
          <div style="font-weight: 700; font-size: 10px; margin-bottom: 4px;">Architectural Highlights:</div>
          <ul class="bullet-list">
            <li>Modern mono-pitch shed roofline with clerestory transoms</li>
            <li>Vaulted Great Room with mezzanine sleeping loft above</li>
            <li>Ground floor private bedroom with large closet & egress</li>
            <li>Board-and-batten dark charcoal siding with natural oak inlays</li>
            <li>Front and rear entry decks with step lighting</li>
          </ul>
        </div>
        <div class="ideal-tag">Ideal Placement: Medium-to-large residential lots, remote tech professionals.</div>
      </div>
    </div>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 7 of 17</span>
    </div>
  </div>

  <!-- SLIDE 8: SIERRA FLOORPLAN -->
  <div class="slide">
    <div class="slide-header">
      <div class="sub-tag">MODEL 03 &bull; ARCHITECTURAL CAD FLOORPLAN</div>
      <h2>The Sierra CAD Blueprint (14' &times; 31'-6" &bull; 420 Sq. Ft.)</h2>
    </div>
    <div class="two-col">
      <div class="fp-img-box">
        <img src="file://{FP_SIERRA}" class="fp-img" />
      </div>
      <div class="info-card">
        <div>
          <div class="badge-dim">SEALED DRAWING NUMBER: A1 &bull; S2A MODULAR</div>
          <div style="font-weight: 800; font-size: 12px; color: #0f172a; margin-bottom: 6px;">Detailed Room Layout:</div>
          <ul class="bullet-list">
            <li><strong>Vaulted Great Room:</strong> 30'-0" long open great room with double glass entry sliders.</li>
            <li><strong>Open to Loft:</strong> Mezzanine sleeping/storage loft overhead with access ladder.</li>
            <li><strong>Private Bedroom:</strong> Ground floor master with large wardrobe and 3660SH egress.</li>
            <li><strong>Gourmet Kitchen:</strong> 4836SL window over sink, cooking range, and pantry.</li>
            <li><strong>Dual Stoops:</strong> Site-built front and rear outdoor entertaining decks.</li>
          </ul>
        </div>
        <div style="font-size: 9px; font-weight: 700; color: #16a34a;">✓ Pre-Approved California SB 9 / Statewide ADU Building Code Compliant</div>
      </div>
    </div>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 8 of 17</span>
    </div>
  </div>

  <!-- SLIDE 9: MEADOW 3D -->
  <div class="slide">
    <div class="slide-header">
      <div class="sub-tag">MODEL 04 &bull; 3D EXTERIOR & REVENUE</div>
      <h2>The Meadow &bull; 15' &times; 31'-9" (435 Sq. Ft.)</h2>
    </div>
    <div class="two-col">
      <img src="file://{MEADOW_IMG}" class="model-img" />
      <div class="info-card">
        <div>
          <div class="badge-dim">DIMENSIONS: 15' &times; 31'-9" &bull; 435 SQ. FT.</div>
          <div class="rent-val">Est. Monthly Rent: $2,600 &ndash; $3,400 / mo</div>
          <div class="split-val">Homeowner Passive Net: $700 &ndash; $1,000 / mo</div>
          <div style="font-weight: 700; font-size: 10px; margin-bottom: 4px;">Architectural Highlights:</div>
          <ul class="bullet-list">
            <li>Full-footprint rooftop sun terrace with slatted privacy railings</li>
            <li>Exterior open staircase connecting yard directly to rooftop</li>
            <li>Symmetrical front facade with flanked oversized black windows</li>
            <li>Dual lofts (sleeping + storage/lounge) and plant ledge top</li>
            <li>Full gourmet kitchen with dining peninsula and rear patio door</li>
          </ul>
        </div>
        <div class="ideal-tag">Ideal Placement: Wide suburban parcels (50'+ width), high-ADR luxury travel rentals.</div>
      </div>
    </div>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 9 of 17</span>
    </div>
  </div>

  <!-- SLIDE 10: MEADOW FLOORPLAN -->
  <div class="slide">
    <div class="slide-header">
      <div class="sub-tag">MODEL 04 &bull; ARCHITECTURAL CAD FLOORPLAN</div>
      <h2>The Meadow CAD Blueprint (15' &times; 31'-9" &bull; 435 Sq. Ft.)</h2>
    </div>
    <div class="two-col">
      <div class="fp-img-box">
        <img src="file://{FP_MEADOW}" class="fp-img" />
      </div>
      <div class="info-card">
        <div>
          <div class="badge-dim">SEALED DRAWING NUMBER: A1 &bull; S2A MODULAR</div>
          <div style="font-weight: 800; font-size: 12px; color: #0f172a; margin-bottom: 6px;">Detailed Room Layout:</div>
          <ul class="bullet-list">
            <li><strong>Dual Lofts:</strong> Two expansive overhead mezzanine lofts plus plant ledge.</li>
            <li><strong>Exterior Staircase:</strong> 2'-9" side exterior open stairs ascending to rooftop deck.</li>
            <li><strong>Rear 6/0 Patio Sliders:</strong> Seamless connection from Great Room to back garden.</li>
            <li><strong>Ground Master Suite:</strong> Private bedroom with closet and stack laundry closet.</li>
            <li><strong>Full 15' Footprint:</strong> S2A's widest modular single-chassis luxury platform.</li>
          </ul>
        </div>
        <div style="font-size: 9px; font-weight: 700; color: #16a34a;">✓ Pre-Approved California SB 9 / Statewide ADU Building Code Compliant</div>
      </div>
    </div>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 10 of 17</span>
    </div>
  </div>

  <!-- SLIDE 11: CASCADE 3D -->
  <div class="slide">
    <div class="slide-header">
      <div class="sub-tag">MODEL 05 &bull; 3D EXTERIOR & REVENUE</div>
      <h2>The Cascade &bull; 12' &times; 27'-9" (300 Sq. Ft.)</h2>
    </div>
    <div class="two-col">
      <img src="file://{CASCADE_IMG}" class="model-img" />
      <div class="info-card">
        <div>
          <div class="badge-dim">DIMENSIONS: 12' &times; 27'-9" &bull; 300 SQ. FT.</div>
          <div class="rent-val">Est. Monthly Rent: $1,650 &ndash; $2,100 / mo</div>
          <div class="split-val">Homeowner Passive Net: $500 &ndash; $750 / mo</div>
          <div style="font-weight: 700; font-size: 10px; margin-bottom: 4px;">Architectural Highlights:</div>
          <ul class="bullet-list">
            <li>Slimline 12-foot exterior profile engineered for narrow lots</li>
            <li>Industrial steel ship's ladder ascending to observation deck</li>
            <li>Efficient open Living/Dining, Kitchenette, and Bath</li>
            <li>Ground floor bedroom with upper loft overhead</li>
            <li>Fits 95% of standard residential parcels without setback variance</li>
          </ul>
        </div>
        <div class="ideal-tag">Ideal Placement: Narrow side yards, urban infill parcels, maximum ROI density.</div>
      </div>
    </div>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 11 of 17</span>
    </div>
  </div>

  <!-- SLIDE 12: CASCADE FLOORPLAN -->
  <div class="slide">
    <div class="slide-header">
      <div class="sub-tag">MODEL 05 &bull; ARCHITECTURAL CAD FLOORPLAN</div>
      <h2>The Cascade CAD Blueprint (12' &times; 27'-9" &bull; 300 Sq. Ft.)</h2>
    </div>
    <div class="two-col">
      <div class="fp-img-box">
        <img src="file://{FP_CASCADE}" class="fp-img" />
      </div>
      <div class="info-card">
        <div>
          <div class="badge-dim">SEALED DRAWING NUMBER: A1 &bull; S2A MODULAR</div>
          <div style="font-weight: 800; font-size: 12px; color: #0f172a; margin-bottom: 6px;">Detailed Room Layout:</div>
          <ul class="bullet-list">
            <li><strong>12' Slim Infill Profile:</strong> Fits 35' narrow residential parcels with 4' setbacks.</li>
            <li><strong>Ship's Ladder:</strong> 2'-9" steel ladder leading to observation rooftop deck.</li>
            <li><strong>Open Living & Kitchenette:</strong> Built-in cabinetry, 4836SL window, and linen storage.</li>
            <li><strong>Ground Bedroom:</strong> Private room with 4/0 Bifold closet and overhead loft.</li>
            <li><strong>Rear Patio Door:</strong> 6/0 sliding glass patio doors for outdoor garden deck.</li>
          </ul>
        </div>
        <div style="font-size: 9px; font-weight: 700; color: #16a34a;">✓ Pre-Approved California SB 9 / Statewide ADU Building Code Compliant</div>
      </div>
    </div>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 12 of 17</span>
    </div>
  </div>

  <!-- SLIDE 13: INTERIOR -->
  <div class="slide">
    <div class="slide-header">
      <div class="sub-tag">INTERIOR ARCHITECTURE & FINISHES</div>
      <h2>Scandinavian Modern Luxury & Space Optimization</h2>
    </div>
    <div class="two-col">
      <img src="file://{INTERIOR_IMG}" class="model-img" />
      <div class="info-card">
        <div>
          <div class="badge-dim">PREMIUM FACTORY-FINISHED INTERIORS</div>
          <div style="font-weight: 700; font-size: 11px; margin-bottom: 6px;">Luxury Standard Package:</div>
          <ul class="bullet-list">
            <li><strong>Vaulted Timber Ceilings:</strong> Exposed natural wood trusses create expansive spatial volume.</li>
            <li><strong>Mezzanine Sleeping Lofts:</strong> Adds 100+ sq ft of functional bonus space.</li>
            <li><strong>Quartz Waterfall Kitchens:</strong> Durable quartz counters, custom oak flat-panel cabinets.</li>
            <li><strong>European Oak Hardwood:</strong> Waterproof luxury engineered plank flooring throughout.</li>
            <li><strong>Designer Bathrooms:</strong> Walk-in glass showers, floating vanities, matte black hardware.</li>
            <li><strong>Energy Efficiency:</strong> All-electric mini-split HVAC, smart thermostats.</li>
          </ul>
        </div>
        <div class="ideal-tag">Zero Compromise Design &bull; Fully Furnished Turnkey Options Available</div>
      </div>
    </div>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 13 of 17</span>
    </div>
  </div>

  <!-- SLIDE 14: REAR PATIO -->
  <div class="slide">
    <div class="slide-header">
      <div class="sub-tag">SEAMLESS INDOOR-OUTDOOR LIVING</div>
      <h2>Private Backyard Entertaining Patios & Sun Decks</h2>
    </div>
    <div class="two-col">
      <img src="file://{PATIO_IMG}" class="model-img" />
      <div class="info-card">
        <div>
          <div class="badge-dim">ELEVATING THE BACKYARD EXPERIENCE</div>
          <div style="font-weight: 700; font-size: 11px; margin-bottom: 6px;">Outdoor Living Highlights:</div>
          <ul class="bullet-list">
            <li><strong>Floor-to-Ceiling Sliders:</strong> Seamless transition between interior and private garden deck.</li>
            <li><strong>Dedicated Outdoor Dining:</strong> Space for outdoor dining tables, lounge chairs, and BBQ.</li>
            <li><strong>Integrated Step Lighting:</strong> Warm low-voltage LED pathway lighting for evening ambiance.</li>
            <li><strong>Privacy Screening:</strong> Architectural slat walls protect homeowner and tenant privacy.</li>
            <li><strong>Minimal Yard Footprint:</strong> Leaves 60-70% of standard backyards open for landscaping.</li>
          </ul>
        </div>
        <div class="ideal-tag">Enhances Main Property Value by 20–30% Upon Installation</div>
      </div>
    </div>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 14 of 17</span>
    </div>
  </div>

  <!-- SLIDE 15: PORTFOLIO MATRIX -->
  <div class="slide">
    <div class="slide-header">
      <div class="sub-tag">PORTFOLIO COMPARISON</div>
      <h2>Comprehensive S2A Modular ADU Lineup</h2>
    </div>
    <table>
      <thead>
        <tr>
          <th>Model</th>
          <th>Dimensions</th>
          <th>Living Area</th>
          <th>Key Blueprint Feature</th>
          <th>Est. Monthly Rent</th>
          <th>Homeowner Net Share</th>
        </tr>
      </thead>
      <tbody>
        <tr>
          <td><strong style="color: #2563eb;">Cascade</strong></td>
          <td>12' &times; 27'-9"</td>
          <td>300 sq. ft.</td>
          <td>12' Slim Infill &bull; Rooftop Observation Ladder</td>
          <td><strong>$1,650 &ndash; $2,100</strong></td>
          <td><strong style="color: #16a34a;">$500 &ndash; $750 / mo</strong></td>
        </tr>
        <tr>
          <td><strong style="color: #2563eb;">Harmony</strong></td>
          <td>14' &times; 24'-0"</td>
          <td>348 sq. ft.</td>
          <td>Pitched Roof &bull; Covered Portico Entry</td>
          <td><strong>$1,850 &ndash; $2,350</strong></td>
          <td><strong style="color: #16a34a;">$550 &ndash; $800 / mo</strong></td>
        </tr>
        <tr>
          <td><strong style="color: #2563eb;">Haven</strong></td>
          <td>14' &times; 29'-0"</td>
          <td>350 sq. ft.</td>
          <td>4' Porch &bull; Rooftop Sky Lounge Deck</td>
          <td><strong>$2,200 &ndash; $2,950</strong></td>
          <td><strong style="color: #16a34a;">$650 &ndash; $950 / mo</strong></td>
        </tr>
        <tr>
          <td><strong style="color: #2563eb;">Sierra</strong></td>
          <td>14' &times; 31'-6"</td>
          <td>420 sq. ft.</td>
          <td>Mono-Pitch Shed &bull; Clerestory Transom Loft</td>
          <td><strong>$2,400 &ndash; $3,100</strong></td>
          <td><strong style="color: #16a34a;">$600 &ndash; $900 / mo</strong></td>
        </tr>
        <tr>
          <td><strong style="color: #2563eb;">Meadow</strong></td>
          <td>15' &times; 31'-9"</td>
          <td>435 sq. ft.</td>
          <td>Dual Lofts &bull; Full Rooftop Terrace Sun Deck</td>
          <td><strong>$2,600 &ndash; $3,400</strong></td>
          <td><strong style="color: #16a34a;">$700 &ndash; $1,000 / mo</strong></td>
        </tr>
      </tbody>
    </table>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 15 of 17</span>
    </div>
  </div>

  <!-- SLIDE 16: TIMELINE -->
  <div class="slide">
    <div class="slide-header">
      <div class="sub-tag">TURNKEY DELIVERY TIMELINE</div>
      <h2>From Site Audit to Monthly Cash Flow in 30–45 Days</h2>
    </div>
    <div style="display: flex; flex-direction: column; gap: 8px; margin-top: 4px;">
      <div class="info-card" style="padding: 8px 12px;">
        <div style="font-size: 11px; font-weight: 800; color: #2563eb;">Day 1–5: Site Feasibility & Digital Permit</div>
        <div style="font-size: 9.5px; color: #64748b;">GIS setback validation, utility hookup mapping, pre-approved architectural plan submission.</div>
      </div>
      <div class="info-card" style="padding: 8px 12px;">
        <div style="font-size: 11px; font-weight: 800; color: #2563eb;">Day 6–20: Offsite Factory Modular Build</div>
        <div style="font-size: 9.5px; color: #64748b;">Precision factory framing, plumbing, electrical, quartz counters, and paint under climate-controlled conditions.</div>
      </div>
      <div class="info-card" style="padding: 8px 12px;">
        <div style="font-size: 11px; font-weight: 800; color: #2563eb;">Day 21–25: Foundation & Trenching</div>
        <div style="font-size: 9.5px; color: #64748b;">Helical pile / concrete pad foundation and underground sewer/power/water utility trenching in yard.</div>
      </div>
      <div class="info-card" style="padding: 8px 12px;">
        <div style="font-size: 11px; font-weight: 800; color: #2563eb;">Day 26–30: Crane Delivery & Installation</div>
        <div style="font-size: 9.5px; color: #64748b;">Single-day crane lift over main house roofline, utility tie-in, and final city inspection sign-off.</div>
      </div>
      <div class="info-card" style="padding: 8px 12px;">
        <div style="font-size: 11px; font-weight: 800; color: #16a34a;">Day 31+: Tenant Placement & Monthly Revenue</div>
        <div style="font-size: 9.5px; color: #64748b;">Furnishing, professional photography, tenant onboarding, and automated monthly revenue-share payouts.</div>
      </div>
    </div>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 16 of 17</span>
    </div>
  </div>

  <!-- SLIDE 17: CTA -->
  <div class="slide slide-dark" style="justify-content: center; text-align: center; padding: 0.8in;">
    <div class="sub-tag" style="margin-bottom: 8px;">UNLOCK YOUR BACKYARD'S PASSIVE INCOME</div>
    <h1 style="font-size: 26px; font-weight: 800; color: #ffffff; margin: 0 0 14px 0;">
      Get Your Free 48-Hour Site Feasibility &<br/>Revenue-Share Audit
    </h1>
    <p style="font-size: 12px; color: #93c5fd; max-width: 7in; margin: 0 auto 20px auto; line-height: 1.5;">
      We assess your parcel geometry, local ADU zoning, utility trenching distances, and provide a guaranteed monthly passive income estimate.
    </p>
    <div style="font-size: 13px; font-weight: 700; color: #4ade80;">
      S2A Modular &bull; Housing-as-a-Service (HaaS) Solutions
    </div>
    <div class="footer-bar">
      <span>S2A Modular &bull; Housing-as-a-Service Portfolio</span>
      <span>Slide 17 of 17</span>
    </div>
  </div>

</body>
</html>"""

    with open(html_file, 'w', encoding='utf-8') as f:
        f.write(html)

    print("Rendering 17-slide landscape PDF via headless Google Chrome...")
    subprocess.run([
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "--headless=new",
        "--disable-gpu",
        "--print-to-pdf-no-header",
        f"--print-to-pdf={PDF_OUTPUT_WORKSPACE}",
        html_file
    ], check=True)

    import shutil
    shutil.copyfile(PDF_OUTPUT_WORKSPACE, PDF_OUTPUT_DOWNLOADS)
    shutil.copyfile(PDF_OUTPUT_WORKSPACE, PDF_OUTPUT_ARTIFACT)
    # Also update the web app public downloads
    web_downloads_pdf = os.path.join(BASE_DIR, "haas-adu-configurator", "public", "downloads", "S2A_Modular_ADU_HaaS_Portfolio.pdf")
    web_downloads_pptx = os.path.join(BASE_DIR, "haas-adu-configurator", "public", "downloads", "S2A_Modular_ADU_HaaS_Portfolio.pptx")
    shutil.copyfile(PDF_OUTPUT_WORKSPACE, web_downloads_pdf)
    shutil.copyfile(PPTX_OUTPUT_WORKSPACE, web_downloads_pptx)
    print(f"PDF successfully created (17 pages):\n- Downloads: {PDF_OUTPUT_DOWNLOADS}\n- Workspace: {PDF_OUTPUT_WORKSPACE}")

if __name__ == "__main__":
    create_pptx()
    create_pdf()
